"""
academic_system/document_processor.py
=======================================
Text extraction and chunking for the Document API (v3.3.0+).

Supported formats
-----------------
  .pdf    — pdfplumber (primary), PyPDF2 (fallback)
  .docx   — python-docx
  .txt    — plain UTF-8 / Latin-1
  .md     — treated as plain text
  .pptx   — python-pptx (slide text extraction)

Install
-------
  pip install pdfplumber python-docx python-pptx PyPDF2

Public API
----------
  extract_text(file_path)          → DocumentResult
  chunk_text(text, chunk_size, overlap) → List[str]

DocumentResult fields
---------------------
  text          str        Full extracted text (may be very long)
  chunks        List[str]  Fixed-size overlapping text windows
  page_count    int        Number of pages / slides (0 for txt/md)
  word_count    int
  char_count    int
  title         str        Inferred from filename or doc metadata
  file_type     str        "pdf" | "docx" | "txt" | "pptx" | "md"
  error         str | None Non-fatal extraction warning
"""
from __future__ import annotations

import logging
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# ── pdfplumber (primary PDF backend) ─────────────────────────────────────────
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
    logger.info("[DocProc] pdfplumber available ✓")
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logger.warning("[DocProc] pdfplumber not installed. pip install pdfplumber")

# ── PyPDF2 (fallback PDF backend) ────────────────────────────────────────────
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
    logger.info("[DocProc] PyPDF2 available ✓ (fallback)")
except ImportError:
    PYPDF2_AVAILABLE = False

# ── python-docx ──────────────────────────────────────────────────────────────
try:
    import docx as _docx
    DOCX_AVAILABLE = True
    logger.info("[DocProc] python-docx available ✓")
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("[DocProc] python-docx not installed. pip install python-docx")

# ── python-pptx ──────────────────────────────────────────────────────────────
try:
    from pptx import Presentation as _PptxPresentation
    PPTX_AVAILABLE = True
    logger.info("[DocProc] python-pptx available ✓")
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("[DocProc] python-pptx not installed. pip install python-pptx")

# ── Supported extensions ──────────────────────────────────────────────────────
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".pptx"}

# ── Default chunking parameters ───────────────────────────────────────────────
_DEFAULT_CHUNK_SIZE    = int(os.environ.get("DOC_CHUNK_SIZE",    "1500"))  # chars
_DEFAULT_CHUNK_OVERLAP = int(os.environ.get("DOC_CHUNK_OVERLAP", "200"))   # chars
_MAX_TOTAL_CHARS       = int(os.environ.get("DOC_MAX_CHARS",     "150000"))# safety cap


@dataclass
class DocumentResult:
    text:       str
    chunks:     List[str]
    page_count: int        = 0
    word_count: int        = 0
    char_count: int        = 0
    title:      str        = ""
    file_type:  str        = ""
    error:      Optional[str] = None
    headings:   List[str]  = field(default_factory=list)  # extracted section headings


# ──────────────────────────────────────────────────────────────────────────────
#  CHUNKING
# ──────────────────────────────────────────────────────────────────────────────

def chunk_text(
    text:       str,
    chunk_size: int = _DEFAULT_CHUNK_SIZE,
    overlap:    int = _DEFAULT_CHUNK_OVERLAP,
) -> List[str]:
    """
    Split text into overlapping fixed-size windows.

    Strategy:
      1. Try to split on paragraph boundaries (double newline) first.
      2. If a paragraph is longer than chunk_size, fall back to hard split.
      3. Overlap is achieved by carrying the last `overlap` chars of the
         previous chunk into the next one, preserving context across boundaries.

    Returns [] if text is empty.
    """
    text = text.strip()
    if not text:
        return []

    # Single chunk — text fits entirely
    if len(text) <= chunk_size:
        return [text]

    # Split on paragraph boundaries
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]

    chunks:  List[str] = []
    current: str       = ""

    for para in paragraphs:
        # Para fits in current chunk
        if len(current) + len(para) + 2 <= chunk_size:
            current = (current + "\n\n" + para).strip() if current else para
        else:
            # Flush current chunk
            if current:
                chunks.append(current)
                # Start new chunk with overlap from previous
                overlap_text = current[-overlap:] if overlap > 0 else ""
                current      = (overlap_text + "\n\n" + para).strip() if overlap_text else para
            else:
                # Para itself is larger than chunk_size — hard split
                for i in range(0, len(para), chunk_size - overlap):
                    piece = para[i : i + chunk_size]
                    if piece.strip():
                        chunks.append(piece.strip())
                current = ""

    if current.strip():
        chunks.append(current.strip())

    return chunks if chunks else [text[:chunk_size]]


# ──────────────────────────────────────────────────────────────────────────────
#  EXTRACTORS
# ──────────────────────────────────────────────────────────────────────────────

def _extract_pdf(file_path: str) -> tuple[str, int, Optional[str]]:
    """Returns (text, page_count, error_message_or_None)."""
    # Primary: pdfplumber
    if PDFPLUMBER_AVAILABLE:
        try:
            pages_text: List[str] = []
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        pages_text.append(t)
            text = "\n\n".join(pages_text)
            if text.strip():
                return text, page_count, None
            # pdfplumber got nothing — fall through to PyPDF2
            logger.warning("[DocProc] pdfplumber returned empty text — trying PyPDF2.")
        except Exception as exc:
            logger.warning(f"[DocProc] pdfplumber failed ({exc}) — trying PyPDF2.")

    # Fallback: PyPDF2
    if PYPDF2_AVAILABLE:
        try:
            pages_text = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                page_count = len(reader.pages)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        pages_text.append(t)
            text = "\n\n".join(pages_text)
            return text, page_count, None
        except Exception as exc:
            return "", 0, f"PyPDF2 failed: {exc}"

    return "", 0, "No PDF extraction library available. pip install pdfplumber"


def _extract_docx(file_path: str) -> tuple[str, int, Optional[str]]:
    if not DOCX_AVAILABLE:
        return "", 0, "python-docx not installed. pip install python-docx"
    try:
        doc        = _docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text       = "\n\n".join(paragraphs)
        return text, len(doc.paragraphs), None
    except Exception as exc:
        return "", 0, f"python-docx failed: {exc}"


def _extract_pptx(file_path: str) -> tuple[str, int, Optional[str]]:
    if not PPTX_AVAILABLE:
        return "", 0, "python-pptx not installed. pip install python-pptx"
    try:
        prs        = _PptxPresentation(file_path)
        slides_text: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
            if slide_lines:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
        text = "\n\n".join(slides_text)
        return text, len(prs.slides), None
    except Exception as exc:
        return "", 0, f"python-pptx failed: {exc}"


def _extract_txt(file_path: str) -> tuple[str, int, Optional[str]]:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            with open(file_path, encoding=enc) as f:
                text = f.read()
            return text, 0, None
        except UnicodeDecodeError:
            continue
    return "", 0, "Could not decode text file — unknown encoding."


def _extract_headings(text: str) -> List[str]:
    """Extract Markdown-style or ALL-CAPS headings from text (best-effort)."""
    headings = []
    for line in text.splitlines():
        line = line.strip()
        # Markdown heading
        if line.startswith("#"):
            headings.append(re.sub(r'^#+\s*', '', line).strip())
        # Short ALL-CAPS line (likely a section header in PDFs)
        elif len(line) > 4 and line.isupper() and len(line.split()) <= 8:
            headings.append(line)
    return headings[:20]  # cap at 20


# ──────────────────────────────────────────────────────────────────────────────
#  PUBLIC API
# ──────────────────────────────────────────────────────────────────────────────

def extract_text(file_path: str) -> DocumentResult:
    """
    Extract text from a document file and return a DocumentResult.

    Handles: .pdf, .docx, .txt, .md, .pptx

    The returned `text` is the full raw extraction (capped at
    _MAX_TOTAL_CHARS to prevent OOM on huge documents).
    The returned `chunks` are ready to pass into LLM prompts.
    """
    path      = Path(file_path)
    ext       = path.suffix.lower()
    file_type = ext.lstrip(".")

    if ext not in DOCUMENT_EXTENSIONS:
        return DocumentResult(
            text="", chunks=[], file_type=file_type,
            title=path.stem,
            error=f"Unsupported file type '{ext}'. Supported: {sorted(DOCUMENT_EXTENSIONS)}",
        )

    # ── Extract ───────────────────────────────────────────────────────────────
    if ext == ".pdf":
        text, page_count, err = _extract_pdf(file_path)
    elif ext == ".docx":
        text, page_count, err = _extract_docx(file_path)
    elif ext == ".pptx":
        text, page_count, err = _extract_pptx(file_path)
    else:  # .txt, .md
        text, page_count, err = _extract_txt(file_path)

    # ── Safety cap ────────────────────────────────────────────────────────────
    if len(text) > _MAX_TOTAL_CHARS:
        logger.warning(
            f"[DocProc] {path.name}: text truncated from {len(text):,} "
            f"to {_MAX_TOTAL_CHARS:,} chars."
        )
        text = text[:_MAX_TOTAL_CHARS]
        err  = (err or "") + f" [Text truncated to {_MAX_TOTAL_CHARS:,} chars]"

    # ── Metadata ──────────────────────────────────────────────────────────────
    word_count = len(text.split()) if text else 0
    char_count = len(text)
    title      = path.stem.replace("_", " ").replace("-", " ").title()
    headings   = _extract_headings(text)

    # ── Chunk ─────────────────────────────────────────────────────────────────
    chunks = chunk_text(text)

    logger.info(
        f"[DocProc] {path.name}: {file_type}, {page_count} pages, "
        f"{word_count:,} words, {len(chunks)} chunks"
        + (f", WARNING: {err}" if err else "")
    )

    return DocumentResult(
        text       = text,
        chunks     = chunks,
        page_count = page_count,
        word_count = word_count,
        char_count = char_count,
        title      = title,
        file_type  = file_type,
        error      = err if err else None,
        headings   = headings,
    )